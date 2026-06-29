"""Render the anywork4me summary as a real .mp4 (Pillow frames + ffmpeg) and an editable .pptx."""
import math, os, subprocess, sys

from PIL import Image, ImageDraw, ImageFont

W, H = 1920, 1080
BG = (67, 56, 202)          # #4338ca indigo-700 (brand purple)
WHITE = (255, 255, 255)
L1 = (199, 210, 254)        # #c7d2fe indigo-200
L2 = (165, 180, 252)        # #a5b4fc indigo-300
AMBER = (239, 159, 39)      # #EF9F27
DARK = (49, 46, 129)        # #312e81 indigo-900
GREY = (95, 94, 90)
CHIPBG = (237, 233, 254)    # #ede9fe violet-100
CHIPTX = (109, 40, 217)     # #6d28d9 violet-700
BUB = (165, 180, 252)       # #a5b4fc indigo-300

F = "C:/Windows/Fonts/segoeui.ttf"
FB = "C:/Windows/Fonts/segoeuib.ttf"
OUT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))


def font(bold, size):
    return ImageFont.truetype(FB if bold else F, size)


def wrap(d, text, fnt, maxw):
    words, lines, cur = text.split(), [], ""
    for w in words:
        t = (cur + " " + w).strip()
        if d.textlength(t, font=fnt) <= maxw:
            cur = t
        else:
            if cur:
                lines.append(cur)
            cur = w
    if cur:
        lines.append(cur)
    return lines


def center(d, cx, y, text, fnt, fill, anchor="mm"):
    d.text((cx, y), text, font=fnt, fill=fill, anchor=anchor)


def block(d, cx, y, lines, fnt, fill, lh):
    for i, ln in enumerate(lines):
        d.text((cx, y + i * lh), ln, font=fnt, fill=fill, anchor="mm")
    return y + len(lines) * lh


def tracked(d, cx, y, text, fnt, fill, track=10):
    widths = [d.textlength(c, font=fnt) for c in text]
    total = sum(widths) + track * (len(text) - 1)
    x = cx - total / 2
    for c, wd in zip(text, widths):
        d.text((x, y), c, font=fnt, fill=fill, anchor="lm")
        x += wd + track


def star(d, cx, cy, r, fill):
    pts = []
    for k in range(10):
        ang = -math.pi / 2 + k * math.pi / 5
        rad = r if k % 2 == 0 else r * 0.42
        pts.append((cx + rad * math.cos(ang), cy + rad * math.sin(ang)))
    d.polygon(pts, fill=fill)


def rrect(d, box, r, fill):
    d.rounded_rectangle(box, radius=r, fill=fill)


CX = W // 2


def s_title(d):
    rrect(d, [CX - 38, 300, CX + 38, 376], 18, WHITE)
    d.text((CX, 338), "a", font=font(True, 64), fill=BG, anchor="mm")
    center(d, CX, 500, "anywork4me", font(True, 132), WHITE)
    center(d, CX, 612, "Find Work  ·  Sell Products  ·  Hire People Fast", font(False, 42), L1)
    center(d, CX, 690, "the world's simplest local marketplace", font(False, 32), L2)


def s_problem(d):
    tracked(d, CX, 360, "THE PROBLEM", font(True, 30), L2, 8)
    block(d, CX, 520, wrap(d, "Finding trusted local help is harder than it should be.",
                           font(True, 78), 1400), font(True, 78), WHITE, 96)
    center(d, CX, 760, "Scattered across WhatsApp groups, referrals and guesswork.",
           font(False, 40), L1)


def s_solution(d):
    tracked(d, CX, 340, "THE SOLUTION", font(True, 30), L2, 8)
    center(d, CX, 450, "One simple marketplace. Anywhere.", font(True, 80), WHITE)
    rrect(d, [CX - 400, 590, CX + 400, 700], 26, WHITE)
    d.ellipse([CX - 360, 630, CX - 332, 658], outline=GREY, width=4)
    d.line([CX - 338, 656, CX - 326, 668], fill=GREY, width=4)
    d.text((CX - 300, 645), "Find a DJ near me", font=font(False, 38), fill=DARK, anchor="lm")
    rrect(d, [CX + 250, 624, CX + 372, 676], 26, CHIPBG)
    d.text((CX + 311, 650), "Accra", font=font(False, 30), fill=CHIPTX, anchor="mm")
    center(d, CX, 800, "Search what you need, nearby — in seconds.", font(False, 40), L1)


def s_customers(d):
    tracked(d, CX, 320, "FOR CUSTOMERS", font(True, 30), L2, 8)
    center(d, CX, 420, "Find & hire with confidence.", font(True, 80), WHITE)
    rrect(d, [CX - 360, 540, CX + 360, 760], 28, WHITE)
    d.ellipse([CX - 320, 580, CX - 240, 660], fill=(79, 70, 229))
    d.text((CX - 280, 620), "M", font=font(True, 44), fill=WHITE, anchor="mm")
    d.text((CX - 210, 598), "By Maali", font=font(True, 40), fill=DARK, anchor="lm")
    d.text((CX - 210, 644), "Mechanic · Accra", font=font(False, 32), fill=GREY, anchor="lm")
    sx = CX - 210
    for k in range(5):
        star(d, sx + k * 50, 712, 19, AMBER)
    d.text((sx + 270, 712), "4.9", font=font(True, 34), fill=DARK, anchor="lm")
    center(d, CX, 860, "Real ratings, safe contact, book on the spot.", font(False, 40), L1)


def s_providers(d):
    tracked(d, CX, 360, "FOR PROVIDERS", font(True, 30), L2, 8)
    center(d, CX, 480, "List free. Get found. Get booked.", font(True, 84), WHITE)
    block(d, CX, 650, wrap(d, "Create a profile in minutes — welder to wedding DJ. "
                           "Free to start, Pro to grow.", font(False, 42), 1300),
          font(False, 42), L1, 60)


def s_messaging(d):
    tracked(d, CX, 320, "MESSAGING & BOOKING", font(True, 30), L2, 8)
    center(d, CX, 420, "Chat & book — right on the site.", font(True, 78), WHITE)
    rrect(d, [CX - 430, 540, CX - 40, 620], 24, WHITE)
    d.text((CX - 410, 580), "Hi, available this weekend?", font=font(False, 32), fill=DARK, anchor="lm")
    rrect(d, [CX + 30, 660, CX + 420, 740], 24, BUB)
    d.text((CX + 50, 700), "Yes — let's book it.", font=font(False, 32), fill=DARK, anchor="lm")
    center(d, CX, 850, "Two-way messages and instant booking alerts, in-app.", font(False, 40), L1)


def s_global(d):
    tracked(d, CX, 340, "BUILT TO SCALE", font(True, 30), L2, 8)
    center(d, CX, 450, "Country-agnostic. Mobile-first.", font(True, 80), WHITE)
    chips = ["GH₵", "₦", "$", "KSh", "R"]
    gap = 26
    fnt = font(True, 40)
    ws = [max(96, d.textlength(c, font=fnt) + 60) for c in chips]
    total = sum(ws) + gap * (len(chips) - 1)
    x = CX - total / 2
    for c, wd in zip(chips, ws):
        rrect(d, [x, 600, x + wd, 684], 42, WHITE)
        d.text((x + wd / 2, 642), c, font=fnt, fill=BG, anchor="mm")
        x += wd + gap
    center(d, CX, 800, "Any country, any currency — one platform.", font(False, 40), L1)


def s_cta(d):
    center(d, CX, 430, "anywork4me.com", font(True, 104), WHITE)
    center(d, CX, 560, "Find Work  ·  Sell Products  ·  Hire People Fast", font(False, 42), L1)
    rrect(d, [CX - 200, 660, CX + 200, 752], 46, WHITE)
    center(d, CX, 706, "Join free today", font(True, 42), BG)


SCENES = [
    ("THE PROBLEM", s_title, "anywork4me", "the world's simplest local marketplace"),
    ("THE PROBLEM", s_problem, "Finding trusted local help is harder than it should be.",
     "Scattered across WhatsApp groups, referrals and guesswork."),
    ("THE SOLUTION", s_solution, "One simple marketplace. Anywhere.",
     "Search what you need, nearby — in seconds."),
    ("FOR CUSTOMERS", s_customers, "Find & hire with confidence.",
     "Real ratings, safe contact, book on the spot."),
    ("FOR PROVIDERS", s_providers, "List free. Get found. Get booked.",
     "Create a profile in minutes — welder to wedding DJ. Free to start, Pro to grow."),
    ("MESSAGING & BOOKING", s_messaging, "Chat & book — right on the site.",
     "Two-way messages and instant booking alerts, in-app."),
    ("BUILT TO SCALE", s_global, "Country-agnostic. Mobile-first.",
     "Any country, any currency — one platform."),
    ("JOIN", s_cta, "anywork4me.com", "Find Work · Sell Products · Hire People Fast"),
]

draws = [s_title, s_problem, s_solution, s_customers, s_providers, s_messaging, s_global, s_cta]
frames = []
for i, fn in enumerate(draws):
    img = Image.new("RGB", (W, H), BG)
    fn(ImageDraw.Draw(img))
    p = os.path.join(OUT, f"_frame_{i}.png")
    img.save(p)
    frames.append(p)
print("frames:", len(frames))

# ---- mp4 via ffmpeg xfade, fallback to hard-cut concat ----
mp4 = os.path.join(OUT, "anywork4me_summary.mp4")
dur, xf, clip = 4.0, 0.6, 4.6
ins = []
for p in frames:
    ins += ["-loop", "1", "-t", str(clip), "-i", p]
parts, prev, off = [], "[0:v]", dur
for i in range(1, len(frames)):
    out = f"[v{i}]" if i < len(frames) - 1 else "[vout]"
    parts.append(f"{prev}[{i}:v]xfade=transition=fade:duration={xf}:offset={off:.2f}{out}")
    prev, off = out, off + dur
fc = ";".join(parts)
cmd = ["ffmpeg", "-y", *ins, "-filter_complex", fc, "-map", "[vout]",
       "-r", "30", "-pix_fmt", "yuv420p", "-c:v", "libx264", "-movflags", "+faststart", mp4]
r = subprocess.run(cmd, capture_output=True, text=True)
if r.returncode != 0:
    print("xfade failed, using hard-cut concat\n", r.stderr[-600:])
    lst = os.path.join(OUT, "_list.txt")
    with open(lst, "w") as f:
        for p in frames:
            f.write(f"file '{p.replace(os.sep, '/')}'\nduration {dur}\n")
        f.write(f"file '{frames[-1].replace(os.sep, '/')}'\n")
    subprocess.run(["ffmpeg", "-y", "-f", "concat", "-safe", "0", "-i", lst,
                    "-r", "30", "-pix_fmt", "yuv420p", "-c:v", "libx264", mp4], check=True)
print("mp4:", os.path.exists(mp4), round(os.path.getsize(mp4) / 1024) if os.path.exists(mp4) else 0, "KB")

# ---- editable pptx ----
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
blank = prs.slide_layouts[6]
for kicker, _, title, sub in SCENES:
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = RGBColor(0x43, 0x38, 0xCA)

    def tb(top, height, text, size, bold, rgb, sz_align=PP_ALIGN.CENTER):
        box = s.shapes.add_textbox(Inches(1), Inches(top), Inches(11.333), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = sz_align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*rgb)
        run.font.name = "Segoe UI"

    tb(2.3, 0.5, kicker, 16, True, (0xA5, 0xB4, 0xFC))
    tb(2.9, 1.8, title, 40, True, (255, 255, 255))
    tb(4.9, 1.0, sub, 20, False, (0xC7, 0xD2, 0xFE))

pptx = os.path.join(OUT, "anywork4me_summary.pptx")
prs.save(pptx)
print("pptx:", os.path.exists(pptx), round(os.path.getsize(pptx) / 1024), "KB")

for p in frames:
    try:
        os.remove(p)
    except OSError:
        pass
print("DONE")
